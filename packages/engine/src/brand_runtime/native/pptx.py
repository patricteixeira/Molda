"""Renderer e inspeção PPTX template-first para o Gate 0 do M2."""

from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.slide import SlideLayout
from pptx.util import Emu, Pt

from brand_runtime.ir.models import BrandIR, SemanticRole
from brand_runtime.kit.models import ContentSpec, LayoutSpec, Slot, TextValue
from brand_runtime.native.ooxml import OoxmlError, validate_ooxml

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_BODY_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.SUBTITLE,
}
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


@dataclass(frozen=True, slots=True)
class SemanticShape:
    """Estado reencontrado no arquivo após edição externa."""

    role: str
    name: str
    kind: str
    text: str | None
    font_family: str | None
    font_size_pt: float | None
    color: str | None


class PptxRenderError(OoxmlError):
    """O template ou os contratos não permitem um PPTX nativo coerente."""


def _compatible_layout(layout: SlideLayout) -> bool:
    placeholder_types = {placeholder.placeholder_format.type for placeholder in layout.placeholders}
    return bool(placeholder_types & _TITLE_TYPES) and bool(placeholder_types & _BODY_TYPES)


def _select_layout(presentation: Presentation, requested_name: str | None) -> SlideLayout:
    if requested_name:
        for layout in presentation.slide_layouts:
            if layout.name == requested_name and _compatible_layout(layout):
                return layout
        raise PptxRenderError(
            f"O layout nativo «{requested_name}» não possui título e corpo editáveis."
        )
    for layout in presentation.slide_layouts:
        if _compatible_layout(layout):
            return layout
    raise PptxRenderError("O template não possui um layout com título e corpo editáveis.")


def _tag_shape(shape, *, role: str, revision_id: str, slot_id: str) -> None:
    name = f"br:{role}:{slot_id}"
    shape.name = name
    nodes = shape._element.xpath(".//p:cNvPr")
    if not nodes:
        raise PptxRenderError(f"Não foi possível etiquetar o shape do slot «{slot_id}».")
    nodes[0].set("name", name)
    nodes[0].set(
        "descr",
        f"brand-role={role};brand-revision={revision_id};slot={slot_id}",
    )


def _set_text(shape, value: TextValue, role: SemanticRole, ir: BrandIR) -> None:
    shape.text = value.text
    paragraph = shape.text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = value.text
    else:
        run = paragraph.runs[0]
    font = ir.fonts[role.font]
    color = ir.colors[role.color]
    run.font.name = font.family
    run.font.size = Pt(role.max_size_px * 0.75)
    run.font.bold = font.weight >= 600
    run.font.italic = font.style == "italic"
    run.font.color.rgb = RGBColor.from_string(color.value.removeprefix("#"))


def _placeholder(slide, accepted: set[PP_PLACEHOLDER]):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type in accepted:
            return placeholder
    return None


def _resolve_asset(path: str, asset_root: Path | None) -> Path:
    candidate = Path(path)
    if not candidate.is_absolute() and asset_root is not None:
        candidate = asset_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise PptxRenderError(f"O asset «{path}» não foi encontrado.")
    try:
        with Image.open(candidate) as image:
            image.verify()
    except (OSError, ValueError) as error:
        raise PptxRenderError(f"O asset «{path}» não é uma imagem raster válida.") from error
    return candidate


def _normalized_box(presentation: Presentation, layout: LayoutSpec, slot: Slot) -> tuple[Emu, ...]:
    x, y, width, height = slot.area
    canvas_width = layout.canvas.width_px
    canvas_height = layout.canvas.height_px
    return (
        Emu(round(presentation.slide_width * x / canvas_width)),
        Emu(round(presentation.slide_height * y / canvas_height)),
        Emu(round(presentation.slide_width * width / canvas_width)),
        Emu(round(presentation.slide_height * height / canvas_height)),
    )


def _position_shape(shape, presentation: Presentation, layout: LayoutSpec, slot: Slot) -> None:
    left, top, width, height = _normalized_box(presentation, layout, slot)
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height


def _text_slots(layout: LayoutSpec, content: ContentSpec) -> list[tuple[Slot, TextValue]]:
    values: list[tuple[Slot, TextValue]] = []
    for slot in layout.slots:
        value = content.values.get(slot.id)
        if slot.kind == "text" and isinstance(value, TextValue):
            values.append((slot, value))
    return values


def _validate_contracts(ir: BrandIR, layout: LayoutSpec, content: ContentSpec) -> None:
    if content.brand_revision_id != ir.revision.id:
        raise PptxRenderError("O Content Spec não pertence à revisão de marca informada.")
    if content.layout_id != layout.id:
        raise PptxRenderError("O Content Spec não pertence ao Layout Spec informado.")
    for slot, value in _text_slots(layout, content):
        if slot.role not in ir.roles:
            raise PptxRenderError(f"O papel semântico «{slot.role}» não existe no Brand IR.")
        role = ir.roles[slot.role]
        if role.font not in ir.fonts or role.color not in ir.colors:
            raise PptxRenderError(f"O papel semântico «{slot.role}» possui referência ausente.")


def render_pptx(
    template_path: Path,
    output_path: Path,
    ir: BrandIR,
    layout: LayoutSpec,
    content: ContentSpec,
    *,
    asset_root: Path | None = None,
    native_layout_name: str | None = None,
) -> Path:
    """Preenche um slide nativo e preserva masters, layouts e theme do template."""
    template_path = template_path.resolve()
    output_path = output_path.resolve()
    if template_path == output_path:
        raise PptxRenderError("O template original nunca pode ser sobrescrito.")
    blocking = [item for item in validate_ooxml(template_path) if item.blocking]
    if blocking:
        raise PptxRenderError(f"O template possui {len(blocking)} erro(s) estrutural(is).")
    _validate_contracts(ir, layout, content)

    presentation = Presentation(template_path)
    native_layout = _select_layout(presentation, native_layout_name)
    slide = presentation.slides.add_slide(native_layout)
    text_values = _text_slots(layout, content)
    if not text_values:
        raise PptxRenderError("O slide precisa de ao menos um slot de texto preenchido.")

    title_slot, title_value = text_values[0]
    title_shape = _placeholder(slide, _TITLE_TYPES)
    if title_shape is None:
        raise PptxRenderError("O layout selecionado perdeu o placeholder de título.")
    _position_shape(title_shape, presentation, layout, title_slot)
    _set_text(title_shape, title_value, ir.roles[title_slot.role], ir)
    _tag_shape(
        title_shape,
        role=title_slot.role,
        revision_id=ir.revision.id,
        slot_id=title_slot.id,
    )

    body_shape = _placeholder(slide, _BODY_TYPES)
    if body_shape is None:
        raise PptxRenderError("O layout selecionado perdeu o placeholder de corpo.")
    if len(text_values) > 1:
        body_slot, body_value = text_values[1]
        _position_shape(body_shape, presentation, layout, body_slot)
        _set_text(body_shape, body_value, ir.roles[body_slot.role], ir)
        _tag_shape(
            body_shape,
            role=body_slot.role,
            revision_id=ir.revision.id,
            slot_id=body_slot.id,
        )
    else:
        body_shape.text = ""
        _tag_shape(body_shape, role="body", revision_id=ir.revision.id, slot_id="body")

    logo_slot = next((slot for slot in layout.slots if slot.kind == "logo"), None)
    if logo_slot is not None:
        asset_token = logo_slot.asset_token or next(iter(ir.assets), None)
        if asset_token is None or asset_token not in ir.assets:
            raise PptxRenderError(
                "O layout exige logo, mas o Brand IR não possui asset compatível."
            )
        asset = _resolve_asset(ir.assets[asset_token].path, asset_root)
        left, top, width, height = _normalized_box(presentation, layout, logo_slot)
        picture = slide.shapes.add_picture(str(asset), left, top, width=width, height=height)
        _tag_shape(
            picture,
            role="logo",
            revision_id=ir.revision.id,
            slot_id=logo_slot.id,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temp_name = tempfile.mkstemp(
        prefix=f".{output_path.stem}-", suffix=".pptx", dir=output_path.parent
    )
    os.close(handle)
    temp_path = Path(temp_name)
    try:
        presentation.save(temp_path)
        blocking = [item for item in validate_ooxml(temp_path) if item.blocking]
        if blocking:
            raise PptxRenderError(f"O PPTX gerado possui {len(blocking)} erro(s) estrutural(is).")
        os.replace(temp_path, output_path)
    finally:
        temp_path.unlink(missing_ok=True)
    return output_path


def _first_text_run(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                return run
    return None


def _theme_style(slide) -> tuple[dict[str, str], dict[str, str]]:
    """Resolve cores e fontes herdadas do tema usado pelo slide."""
    master = slide.slide_layout.slide_master
    theme_part = next(
        (
            relationship.target_part
            for relationship in master.part.rels.values()
            if relationship.reltype.endswith("/theme")
        ),
        None,
    )
    if theme_part is None:
        return {}, {}

    parser = etree.XMLParser(resolve_entities=False, no_network=True, load_dtd=False)
    theme_root = etree.fromstring(theme_part.blob, parser=parser)
    scheme = theme_root.find(
        ".//a:themeElements/a:clrScheme",
        {"a": _A_NS},
    )
    colors: dict[str, str] = {}
    if scheme is not None:
        for slot in scheme:
            color_nodes = list(slot)
            if not color_nodes:
                continue
            color_node = color_nodes[0]
            color_kind = etree.QName(color_node).localname
            value = (
                color_node.get("val")
                if color_kind == "srgbClr"
                else color_node.get("lastClr")
                if color_kind == "sysClr"
                else None
            )
            if value is not None and len(value) == 6:
                colors[etree.QName(slot).localname] = f"#{value.upper()}"

    master_root = etree.fromstring(master.part.blob, parser=parser)
    color_map = master_root.find(".//p:clrMap", {"p": _P_NS})
    if color_map is not None:
        for alias, target in color_map.attrib.items():
            resolved = colors.get(target)
            if resolved is not None:
                colors[alias] = resolved

    fonts: dict[str, str] = {}
    for role, xpath in {
        "major": ".//a:fontScheme/a:majorFont/a:latin",
        "minor": ".//a:fontScheme/a:minorFont/a:latin",
    }.items():
        font = theme_root.find(xpath, {"a": _A_NS})
        typeface = font.get("typeface") if font is not None else None
        if typeface:
            fonts[role] = typeface
    return colors, fonts


def _run_color(run, theme_colors: dict[str, str]) -> str | None:
    color_format = run.font.color
    if color_format.type == MSO_COLOR_TYPE.RGB:
        return f"#{color_format.rgb}"
    if color_format.type == MSO_COLOR_TYPE.SCHEME:
        color_node = getattr(color_format._color, "_xClr", None)
        scheme_name = color_node.get("val") if color_node is not None else None
        if scheme_name is None:
            return None
        return theme_colors.get(scheme_name, f"theme:{scheme_name}")
    return None


def _run_font_family(run, role: str, theme_fonts: dict[str, str]) -> str | None:
    family = run.font.name
    if family not in {None, "", "+mj-lt", "+mn-lt"}:
        return family
    theme_role = "major" if family == "+mj-lt" or role == "heading" else "minor"
    return theme_fonts.get(theme_role)


def _shape_role(shape) -> str | None:
    if shape.name.startswith("br:"):
        parts = shape.name.split(":", 2)
        if len(parts) == 3 and parts[1]:
            return parts[1]
    nodes = shape._element.xpath(".//p:cNvPr")
    if nodes:
        description = nodes[0].get("descr", "")
        for field in description.split(";"):
            key, separator, value = field.partition("=")
            if separator and key == "brand-role" and value:
                return value
    # LibreOffice deliberately rewrites the non-visual properties of native
    # placeholders. The placeholder kind is the third, editor-owned semantic
    # signal and survives that save cycle even when name/descr do not.
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in _TITLE_TYPES:
            return "heading"
        if placeholder_type in _BODY_TYPES:
            return "body"
    return None


def inspect_semantic_shapes(path: Path) -> list[SemanticShape]:
    """Reencontra roles e formatação depois de um save externo ou automatizado."""
    if path.suffix.lower() != ".pptx" or not path.is_file():
        raise PptxRenderError("Informe um arquivo PPTX existente.")
    presentation = Presentation(path)
    findings: list[SemanticShape] = []
    for slide in presentation.slides:
        theme_colors, theme_fonts = _theme_style(slide)
        for shape in slide.shapes:
            role = _shape_role(shape)
            if role is None:
                continue
            run = _first_text_run(shape)
            color = _run_color(run, theme_colors) if run is not None else None
            findings.append(
                SemanticShape(
                    role=role,
                    name=shape.name,
                    kind=("picture" if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else "text"),
                    text=(shape.text if getattr(shape, "has_text_frame", False) else None),
                    font_family=(
                        _run_font_family(run, role, theme_fonts) if run is not None else None
                    ),
                    font_size_pt=(run.font.size.pt if run is not None and run.font.size else None),
                    color=color,
                )
            )
    return findings
