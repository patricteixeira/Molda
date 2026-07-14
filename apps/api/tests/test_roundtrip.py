from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor

from brand_api.models import Job
from brand_api.worker import run_next_job
from brand_runtime import parse_pptx_document_graph, validate_ooxml


def _document(client, compiled):
    response = client.post(
        "/v1/documents",
        json={
            "layoutId": "statement-post-1x1",
            "brandRevisionId": compiled["brandRevisionId"],
            "values": {"headline": {"kind": "text", "text": "Texto original"}},
        },
    )
    assert response.status_code == 201, response.text
    return response.json()["documentId"]


def _run_one(client):
    state = client.app.state
    return run_next_job(
        state.session_factory,
        storage=state.storage,
        exporter=state.exporter,
        settings=state.settings,
    )


def _export_pptx(client, compiled) -> tuple[str, bytes]:
    document_id = _document(client, compiled)
    response = client.post(f"/v1/documents/{document_id}/exports", json={"format": "pptx"})
    assert response.status_code == 202, response.text
    job_id = response.json()["jobId"]
    assert _run_one(client)
    job = client.get(f"/v1/jobs/{job_id}").json()
    assert job["status"] == "succeeded", job
    return job_id, client.get(job["result"]["url"]).content


def _google_style_edit(source: bytes) -> bytes:
    presentation = Presentation(BytesIO(source))
    for index, shape in enumerate(presentation.slides[0].shapes, start=1):
        if not shape.name.startswith("br:"):
            continue
        if shape.name.startswith("br:heading:"):
            shape.text = "Texto aprovado no Google"
            run = shape.text_frame.paragraphs[0].runs[0]
            run.font.color.rgb = RGBColor(0xE5, 0x79, 0x00)
        if shape.name.startswith("br:logo:"):
            shape.left += shape.width
            shape.width *= 2
        shape.name = f"Google Shape;{index};p13"
    edited = BytesIO()
    presentation.save(edited)
    return edited.getvalue()


def test_roundtrip_api_worker_fix_e_storage(client, compiled, tmp_path):
    export_job_id, baseline = _export_pptx(client, compiled)
    edited = _google_style_edit(baseline)

    response = client.post(
        f"/v1/jobs/{export_job_id}/roundtrips",
        files={"file": ("editado.pptx", edited, "application/octet-stream")},
    )
    assert response.status_code == 202, response.text
    lint_job_id = response.json()["jobId"]
    with client.app.state.session_factory() as session:
        queued = session.get(Job, lint_job_id)
        assert queued is not None
        assert queued.kind == "roundtrip-lint"
        edited_sha256 = queued.params["editedSha256"]
    assert client.app.state.storage.get(edited_sha256) == edited

    assert _run_one(client)
    lint_job = client.get(f"/v1/jobs/{lint_job_id}").json()
    assert lint_job["status"] == "succeeded", lint_job
    result = lint_job["result"]
    assert result["kind"] == "roundtrip-lint"
    assert result["documentGraph"]["nodes"][0]["text"] == "Texto aprovado no Google"
    assert result["report"]["summary"]["fixable"] >= 2
    properties = {operation["property"] for operation in result["fixPlan"]["operations"]}
    assert {"color", "boundsPt"} <= properties

    response = client.post(f"/v1/jobs/{lint_job_id}/fixes")
    assert response.status_code == 202, response.text
    fix_job_id = response.json()["jobId"]
    assert _run_one(client)
    fix_job = client.get(f"/v1/jobs/{fix_job_id}").json()
    assert fix_job["status"] == "succeeded", fix_job["error"]
    fixed_result = fix_job["result"]
    assert fixed_result["kind"] == "roundtrip-fix"
    assert fixed_result["fixResult"]["report"]["summary"] == {
        "status": "review",
        "total": 1,
        "info": 1,
        "warning": 0,
        "error": 0,
        "locked": 0,
        "fixable": 0,
    }
    downloaded = client.get(fixed_result["url"])
    assert downloaded.status_code == 200
    assert downloaded.headers["content-type"].endswith("presentationml.presentation")

    corrected = tmp_path / "corrected.pptx"
    corrected.write_bytes(downloaded.content)
    assert not [item for item in validate_ooxml(corrected) if item.blocking]
    graph = parse_pptx_document_graph(corrected)
    assert graph.nodes[0].text == "Texto aprovado no Google"
    assert client.app.state.storage.get(edited_sha256) == edited
    assert not list(client.app.state.settings.work_dir.iterdir())


def test_roundtrip_recusa_upload_invalido_estado_e_auth(client, anon, compiled):
    export_job_id, baseline = _export_pptx(client, compiled)

    invalid = client.post(
        f"/v1/jobs/{export_job_id}/roundtrips",
        files={"file": ("falso.pptx", b"not-a-zip", "application/octet-stream")},
    )
    assert invalid.status_code == 400
    assert invalid.json()["detail"] == "Envie uma apresentação PPTX válida."

    missing = client.post(
        "/v1/jobs/job_000000000000/roundtrips",
        files={"file": ("editado.pptx", baseline, "application/octet-stream")},
    )
    assert missing.status_code == 404

    unauthorized = anon.post(
        f"/v1/jobs/{export_job_id}/roundtrips",
        files={"file": ("editado.pptx", baseline, "application/octet-stream")},
    )
    assert unauthorized.status_code == 401

    unchanged = client.post(
        f"/v1/jobs/{export_job_id}/roundtrips",
        files={"file": ("sem-alteracao.pptx", baseline, "application/octet-stream")},
    )
    assert unchanged.status_code == 202
    unchanged_job_id = unchanged.json()["jobId"]
    assert _run_one(client)
    unchanged_job = client.get(f"/v1/jobs/{unchanged_job_id}").json()
    assert unchanged_job["result"]["report"]["summary"]["status"] == "pass"
    no_fix = client.post(f"/v1/jobs/{unchanged_job_id}/fixes")
    assert no_fix.status_code == 409
    assert no_fix.json()["detail"] == "O arquivo não possui correções automáticas pendentes."
